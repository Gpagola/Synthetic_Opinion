"""Motor de encuestas cuantitativas sobre la población sintética.

- answer_survey: cada persona responde el cuestionario completo en su personaje
  (una llamada al LLM por persona, salida JSON validada contra las opciones).
- compute_results: distribución por pregunta (+ media/NPS) y cruce por segmento.
- export_xlsx: respuestas crudas + tablas por pregunta.
"""

from __future__ import annotations

import io
import json
from concurrent.futures import ThreadPoolExecutor
from types import SimpleNamespace

from openpyxl import Workbook
from sqlalchemy.orm import Session

from app.countries import cultural_context_block
from app.database import SessionLocal
from app.models import Persona, Survey, SurveyResponse
from app.services.llm import get_llm

# Variables de cruce admitidas -> campo del snapshot
BREAK_FIELDS = {
    "genero": "genero", "region": "region", "ingresos": "ingresos",
    "educacion": "nivel_educativo", "edad": "edad",
}
AGE_BANDS = [("18-24", 18, 24), ("25-34", 25, 34), ("35-44", 35, 44), ("45-54", 45, 54),
             ("55-64", 55, 64), ("65-74", 65, 74), ("75-84", 75, 84), ("85+", 85, 200)]
NPS_GROUPS = ["Detractores (0-6)", "Pasivos (7-8)", "Promotores (9-10)"]
LIKERT = ["1", "2", "3", "4", "5"]


def _age_band(edad) -> str:
    if edad is None:
        return "—"
    for label, lo, hi in AGE_BANDS:
        if lo <= edad <= hi:
            return label
    return "85+" if isinstance(edad, int) and edad > 85 else "—"


def _snapshot(p: Persona) -> dict:
    sd = p.sociodemografico or {}
    return {
        "edad": sd.get("edad"), "genero": sd.get("genero"), "region": sd.get("region"),
        "ingresos": sd.get("ingresos"), "nivel_educativo": sd.get("nivel_educativo"),
        "pais_origen": sd.get("pais_origen"),
    }


def _perfil(p: Persona) -> str:
    partes = [f"Bio: {p.bio}" if p.bio else ""]
    if p.sociodemografico:
        partes.append(f"Sociodemografía: {p.sociodemografico}")
    if p.consumidor:
        partes.append(f"Consumo: {p.consumidor}")
    if p.opinion:
        partes.append(f"Valores y opiniones: {p.opinion}")
    return "\n".join(x for x in partes if x)


def _q_spec(q) -> str:
    if q.tipo == "single":
        return f'elige UNA de estas opciones (texto exacto): {q.opciones}'
    if q.tipo == "multiple":
        return f'elige una o varias (array de textos exactos) de: {q.opciones}'
    if q.tipo == "yesno":
        return 'responde "Sí" o "No"'
    if q.tipo == "likert":
        return 'entero del 1 (muy en desacuerdo) al 5 (muy de acuerdo)'
    if q.tipo == "nps":
        return 'entero del 0 al 10 (¿qué probabilidad de recomendarlo?)'
    if q.tipo == "abierta":
        return 'responde con texto libre (1-2 frases) en tu propia voz, como string'
    return ""


def _compute_path(questions: list, respuestas: dict) -> list[int]:
    """Simula el recorrido de una persona por el cuestionario dado sus respuestas.
    Devuelve lista de `orden` (0-based) de las preguntas que habría alcanzado."""
    by_orden = {q.orden: q for q in questions}
    if not by_orden:
        return []
    current = 0
    path = []
    max_steps = len(questions) + 2  # guarda infinita
    steps = 0
    while current is not None and steps < max_steps:
        steps += 1
        q = by_orden.get(current)
        if q is None:
            break
        path.append(current)
        # Buscar salto condicional
        conds = q.condiciones if hasattr(q, "condiciones") else []
        respuesta = respuestas.get(str(q.id), respuestas.get(q.id))
        resp_str = str(respuesta) if respuesta is not None else None
        jumped = False
        for rule in (conds or []):
            si_resp = rule.get("si_respuesta") if isinstance(rule, dict) else getattr(rule, "si_respuesta", None)
            ir_a = rule.get("ir_a_orden") if isinstance(rule, dict) else getattr(rule, "ir_a_orden", None)
            if resp_str == str(si_resp):
                current = ir_a  # None → fin
                jumped = True
                break
        if not jumped:
            # Sin salto: siguiente pregunta en orden
            next_ordenes = sorted(o for o in by_orden if o > current)
            current = next_ordenes[0] if next_ordenes else None
    return path


def _build_prompt(survey: Survey) -> str:
    """Construye el prompt para la IA, incluyendo instrucciones de skip logic si hay condiciones."""
    has_conditions = any(
        (q.condiciones if hasattr(q, "condiciones") else [])
        for q in survey.questions
    )
    lines = []
    for q in survey.questions:
        n = q.orden + 1
        conds = q.condiciones if hasattr(q, "condiciones") else []
        line = f'- P{n} [id={q.id}, orden={q.orden}, tipo={q.tipo}] "{q.texto}" -> {_q_spec(q)}'
        if conds:
            for rule in conds:
                si = rule.get("si_respuesta") if isinstance(rule, dict) else getattr(rule, "si_respuesta", None)
                ir = rule.get("ir_a_orden") if isinstance(rule, dict) else getattr(rule, "ir_a_orden", None)
                destino = f"P{ir + 1}" if ir is not None else "FIN"
                line += f'\n    → Si responde "{si}" → saltar a {destino}'
        lines.append(line)

    instruccion = (
        "Sigue las instrucciones de salto (→): si una condición aplica, salta a la pregunta "
        "indicada y NO respondas las preguntas intermedias. Responde SOLO las preguntas que "
        "alcanzarías. "
        if has_conditions else
        "Responde a TODAS las preguntas. "
    )
    return (
        f"Tema de la encuesta: {survey.tema}\n"
        f"{instruccion}Hazlo de forma honesta y coherente con TODAS tus características "
        f"(edad, ingresos, educación, región, valores, experiencias vitales), no como una "
        f"persona genérica o promedio.\n\n"
        "IMPORTANTE — NO caigas en el sesgo de respuesta central:\n"
        "- En escalas Likert (1-5): usa el 1 o el 5 cuando tu perfil realmente lo justifica. "
        "No te quedes siempre en el 3. Si eres escéptico o muy favorable, dilo con claridad.\n"
        "- En NPS (0-10): los extremos (0-2 o 9-10) son válidos y frecuentes en personas reales. "
        "Úsalos si tu actitud es fuerte. El 5-7 es para indiferencia genuina, no comodidad.\n"
        "- En opciones únicas o múltiples: elige según tu historia y valores, aunque la opción "
        "parezca minoritaria o incómoda. Las personas reales no siempre eligen lo 'razonable'.\n"
        "- En preguntas abiertas: habla con tu voz propia, con tus sesgos y contradicciones. "
        "No des respuestas 'correctas' ni políticamente neutras si tu perfil no lo es.\n\n"
        "Preguntas:\n" + "\n".join(lines) + "\n\n"
        'Devuelve SOLO un JSON {"respuestas": {"<orden>": valor, ...}} usando el ORDEN (0-based) '
        "como clave, con el valor del tipo indicado para cada pregunta respondida. "
        "No añadas texto fuera del JSON."
    )


def _validate(q, raw):
    if raw is None:
        return None
    if q.tipo == "single":
        return raw if raw in q.opciones else None
    if q.tipo == "multiple":
        if isinstance(raw, list):
            vals = [x for x in raw if x in q.opciones]
            return vals or None
        return [raw] if raw in q.opciones else None
    if q.tipo == "yesno":
        s = str(raw).strip().lower()
        if s.startswith("s"):
            return "Sí"
        if s.startswith("n"):
            return "No"
        return None
    if q.tipo in ("likert", "nps"):
        try:
            v = int(round(float(raw)))
        except (ValueError, TypeError):
            return None
        return max(1, min(5, v)) if q.tipo == "likert" else max(0, min(10, v))
    if q.tipo == "abierta":
        s = str(raw).strip()
        return s or None
    return None


def answer_survey(survey_id: int, persona_ids: list[int], modelo: str,
                  reasoning_effort: str | None) -> None:
    db: Session = SessionLocal()
    try:
        survey = db.get(Survey, survey_id)
        if survey is None:
            return
        survey.estado = "running"
        survey.error_msg = None
        survey.modelo = modelo
        survey.reasoning_effort = reasoning_effort
        # limpia respuestas previas (re-lanzamiento)
        db.query(SurveyResponse).filter(SurveyResponse.survey_id == survey_id).delete()
        db.commit()

        personas = {
            p.id: p for p in db.query(Persona).filter(Persona.id.in_(persona_ids)).all()
        }
        ordered = [personas[i] for i in persona_ids if i in personas]
        idioma = survey.idioma
        contexto_pais = cultural_context_block(survey.pais)
        prompt = _build_prompt(survey)
        # Specs de preguntas como objetos planos para validar sin tocar la sesión.
        questions = [SimpleNamespace(
            id=q.id, tipo=q.tipo, opciones=list(q.opciones), orden=q.orden,
            condiciones=list(q.condiciones) if q.condiciones else [],
        ) for q in survey.questions]
        # IMPORTANTE: materializamos todos los datos de persona AQUÍ (hilo principal).
        # Tras los commits, los objetos ORM se expiran; si los hilos accedieran a sus
        # atributos dispararían lazy-loads concurrentes sobre la misma conexión MySQL
        # ("Packet sequence number wrong"). Los hilos solo reciben datos planos.
        items = [(p.id, p.nombre, _perfil(p), _snapshot(p)) for p in ordered]
        # Si el modelo es Claude, usar el proveedor Anthropic; si no, OpenAI.
        llm = get_llm("anthropic" if modelo.startswith("claude") else "openai")

        by_orden = {q.orden: q for q in questions}

        # Para modelos Claude usamos effort="low": las encuestas no necesitan
        # razonamiento profundo y "high" (el default) hace que cada llamada tarde
        # decenas de segundos con 20+ preguntas.
        claude_effort = "low" if modelo.startswith("claude") else None

        def ask(item):
            pid, nombre, perfil, _snap = item
            system = (
                f"Eres {nombre}, una persona real que participa en una encuesta. Respondes en tu "
                f"personaje, en {idioma}. Responde SIEMPRE en JSON válido.\n\n"
                f"Eres un individuo concreto con tu propia historia, sesgos y opiniones. "
                f"Tus respuestas nacen de TODAS tus características a la vez (situación económica, "
                f"educación, región, experiencias vitales, valores), no de una sola. "
                f"No suavices ni promedies: si tu perfil implica desconfianza, frustración, "
                f"entusiasmo o posiciones extremas, exprésalas con claridad.\n\n"
                f"{contexto_pais}\n\nTu perfil:\n{perfil}"
            )
            try:
                data = llm.complete_json(system, prompt, model=modelo,
                                         reasoning_effort=claude_effort or reasoning_effort)
                return pid, data.get("respuestas", {})
            except Exception:  # noqa: BLE001
                return pid, None

        snapshots = {pid: snap for pid, _n, _p, snap in items}
        with ThreadPoolExecutor(max_workers=8) as ex:
            for i, (pid, raw_by_orden) in enumerate(ex.map(ask, items), 1):
                if raw_by_orden is None:
                    continue
                # Construir respuestas normalizadas (clave = str(q.id))
                resp_by_id: dict = {}
                for q in questions:
                    raw_val = raw_by_orden.get(str(q.orden), raw_by_orden.get(q.orden))
                    v = _validate(q, raw_val)
                    if v is not None:
                        resp_by_id[str(q.id)] = v
                # Calcular el path real y conservar solo las preguntas del camino
                path = _compute_path(questions, resp_by_id)
                path_ids = {str(by_orden[o].id) for o in path if o in by_orden}
                resp = {k: v for k, v in resp_by_id.items() if k in path_ids}
                db.add(SurveyResponse(
                    survey_id=survey_id, persona_id=pid,
                    snapshot=snapshots.get(pid), respuestas=resp,
                ))
                if i % 20 == 0:
                    db.commit()
        db.commit()
        survey.estado = "completed"
        db.commit()
    except Exception as exc:  # noqa: BLE001
        db.rollback()
        s = db.get(Survey, survey_id)
        if s is not None:
            s.estado = "error"
            s.error_msg = str(exc)[:1000]
            db.commit()
    finally:
        db.close()


# ---------- Resultados ----------

def _seg_value(snapshot: dict, break_var: str) -> str:
    field = BREAK_FIELDS[break_var]
    val = snapshot.get(field)
    if break_var == "edad":
        return _age_band(val)
    return str(val) if val not in (None, "") else "—"


def _options_for(q) -> list[str]:
    if q.tipo == "yesno":
        return ["Sí", "No"]
    if q.tipo == "likert":
        return LIKERT
    if q.tipo == "nps":
        return NPS_GROUPS
    return list(q.opciones)


def _bucket(q, value):
    """Devuelve la(s) categoría(s) de una respuesta para la distribución."""
    if value is None:
        return []
    if q.tipo == "multiple":
        return list(value)
    if q.tipo == "likert":
        return [str(value)]
    if q.tipo == "nps":
        v = int(value)
        return [NPS_GROUPS[0] if v <= 6 else NPS_GROUPS[1] if v <= 8 else NPS_GROUPS[2]]
    return [str(value)]


def _distrib(q, values: list) -> tuple[list[dict], int]:
    opts = _options_for(q)
    counts = {o: 0 for o in opts}
    n = 0
    for v in values:
        cats = _bucket(q, v)
        if cats:
            n += 1
        for c in cats:
            counts[c] = counts.get(c, 0) + 1
    base = n or 1
    dist = [{"opcion": o, "n": counts.get(o, 0), "pct": round(counts.get(o, 0) * 100 / base, 1)}
            for o in opts]
    return dist, n


def compute_results(db: Session, survey: Survey, break_var: str | None) -> dict:
    responses = db.query(SurveyResponse).filter(SurveyResponse.survey_id == survey.id).all()
    preguntas = []
    for q in survey.questions:
        vals = [r.respuestas.get(str(q.id)) for r in responses if r.respuestas.get(str(q.id)) is not None]
        # Las preguntas abiertas no tienen distribución: se devuelven los verbatims.
        if q.tipo == "abierta":
            textos = [str(v) for v in vals]
            preguntas.append({
                "question_id": q.id, "texto": q.texto, "tipo": q.tipo, "n": len(textos),
                "distribucion": [], "media": None, "nps": None, "cruce": {}, "textos": textos,
            })
            continue
        dist, n = _distrib(q, vals)
        media = None
        nps = None
        if q.tipo in ("likert", "nps"):
            nums = [int(v) for v in vals if isinstance(v, (int, float))]
            if nums:
                media = round(sum(nums) / len(nums), 2)
            if q.tipo == "nps" and nums:
                promo = sum(1 for v in nums if v >= 9)
                detr = sum(1 for v in nums if v <= 6)
                nps = round((promo - detr) * 100 / len(nums), 1)
        cruce = {}
        if break_var in BREAK_FIELDS:
            segs: dict[str, list] = {}
            for r in responses:
                v = r.respuestas.get(str(q.id))
                if v is None:
                    continue
                segs.setdefault(_seg_value(r.snapshot or {}, break_var), []).append(v)
            for seg, sv in sorted(segs.items()):
                cruce[seg] = _distrib(q, sv)[0]
        preguntas.append({
            "question_id": q.id, "texto": q.texto, "tipo": q.tipo, "n": n,
            "distribucion": dist, "media": media, "nps": nps, "cruce": cruce, "textos": [],
        })
    return {
        "estado": survey.estado, "total_respuestas": len(responses),
        "break_var": break_var if break_var in BREAK_FIELDS else None, "preguntas": preguntas,
    }


def export_xlsx(db: Session, survey: Survey) -> bytes:
    responses = db.query(SurveyResponse).filter(SurveyResponse.survey_id == survey.id).all()
    wb = Workbook()

    ws = wb.active
    ws.title = "Respuestas"
    qcols = list(survey.questions)
    ws.append(["persona_id", "edad", "genero", "region", "ingresos", "educacion"]
              + [f"P{i+1}" for i in range(len(qcols))])
    for r in responses:
        s = r.snapshot or {}
        row = [r.persona_id, s.get("edad"), s.get("genero"), s.get("region"),
               s.get("ingresos"), s.get("nivel_educativo")]
        for q in qcols:
            v = r.respuestas.get(str(q.id))
            row.append(", ".join(v) if isinstance(v, list) else v)
        ws.append(row)

    res = compute_results(db, survey, None)
    ws2 = wb.create_sheet("Resultados")
    for q in res["preguntas"]:
        ws2.append([q["texto"]])
        if q["tipo"] == "abierta":
            ws2.append([f"Respuestas abiertas (n={q['n']})"])
            for t in q.get("textos", []):
                ws2.append([t])
            ws2.append([])
            continue
        ws2.append(["Opción", "N", "%"])
        for o in q["distribucion"]:
            ws2.append([o["opcion"], o["n"], o["pct"]])
        if q["media"] is not None:
            ws2.append(["media", q["media"]])
        if q["nps"] is not None:
            ws2.append(["NPS", q["nps"]])
        ws2.append([])

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def export_survey_docx(survey: Survey) -> bytes:
    """Genera un .docx del cuestionario con identidad corporativa Andersen Consulting."""
    from docx import Document
    from docx.shared import Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    NAVY   = RGBColor(0x17, 0x2D, 0x42)   # Heading 1
    BLUE2  = RGBColor(0x21, 0x69, 0x93)   # Tabla header / acento
    BLACK  = RGBColor(0x00, 0x00, 0x00)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    FONT   = "Arial"  # fallback de Helvetica Neue

    def add_h1_border(para):
        """Añade borde inferior a un párrafo (estilo Heading 1 Andersen)."""
        pPr = para._p.get_or_add_pPr()
        pBdr = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), "6")
        bottom.set(qn("w:space"), "1")
        bottom.set(qn("w:color"), "172D42")
        pBdr.append(bottom)
        pPr.append(pBdr)

    def set_cell_bg(cell, hex_color: str):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), hex_color)
        tcPr.append(shd)

    doc = Document()

    # Página A4 con márgenes Andersen
    sec = doc.sections[0]
    sec.page_width  = Cm(21)
    sec.page_height = Cm(29.7)
    sec.top_margin    = Cm(3.3)
    sec.right_margin  = Cm(2)
    sec.bottom_margin = Cm(2.5)
    sec.left_margin   = Cm(2)

    # Footer corporativo
    footer_para = sec.footer.paragraphs[0]
    footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = footer_para.add_run("BRAINTRUST SL es una firma miembro de ANDERSEN CONSULTING")
    run.font.name = FONT
    run.font.size = Pt(8)
    run.font.color.rgb = BLACK

    # ── Título principal ──────────────────────────────────────────
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = title_para.add_run(survey.nombre.upper())
    run.font.name = FONT
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = NAVY
    add_h1_border(title_para)
    doc.add_paragraph()  # espacio

    if survey.tema:
        p = doc.add_paragraph()
        run = p.add_run(survey.tema)
        run.font.name = FONT
        run.font.size = Pt(11)
        run.font.italic = True
        run.font.color.rgb = NAVY

    # ── Introducción ──────────────────────────────────────────────
    if survey.descripcion:
        doc.add_paragraph()
        h = doc.add_paragraph()
        r = h.add_run("INTRODUCCIÓN")
        r.font.name = FONT; r.font.size = Pt(11)
        r.font.bold = True; r.font.color.rgb = NAVY
        add_h1_border(h)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        run = p.add_run(survey.descripcion)
        run.font.name = FONT; run.font.size = Pt(10.5)

    # ── Sección: Cuestionario ─────────────────────────────────────
    doc.add_paragraph()
    h = doc.add_paragraph()
    r = h.add_run("CUESTIONARIO")
    r.font.name = FONT; r.font.size = Pt(11)
    r.font.bold = True; r.font.color.rgb = NAVY
    add_h1_border(h)
    doc.add_paragraph()

    Q_LABELS = {
        "single": "Opción única", "multiple": "Opción múltiple",
        "yesno": "Sí / No", "likert": "Escala 1–5",
        "nps": "NPS 0–10", "abierta": "Respuesta abierta",
    }

    for i, q in enumerate(survey.questions):
        # Número + tipo
        meta = doc.add_paragraph()
        r1 = meta.add_run(f"P{i + 1}. ")
        r1.font.name = FONT; r1.font.size = Pt(10.5)
        r1.font.bold = True; r1.font.color.rgb = BLUE2
        r2 = meta.add_run(f"[{Q_LABELS.get(q.tipo, q.tipo)}]")
        r2.font.name = FONT; r2.font.size = Pt(9)
        r2.font.color.rgb = BLUE2

        # Texto de la pregunta
        texto_p = doc.add_paragraph()
        texto_p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        run = texto_p.add_run(q.texto or "(Sin texto)")
        run.font.name = FONT; run.font.size = Pt(10.5); run.font.bold = True

        # Opciones
        opts = list(q.opciones) if q.opciones else []
        if q.tipo == "yesno":   opts = ["Sí", "No"]
        if q.tipo == "likert":  opts = ["1 – Muy en desacuerdo", "2", "3", "4", "5 – Muy de acuerdo"]
        if q.tipo == "nps":     opts = ["0 ─ 1 ─ 2 ─ 3 ─ 4 ─ 5 ─ 6 ─ 7 ─ 8 ─ 9 ─ 10"]
        if q.tipo == "abierta":
            line = doc.add_paragraph()
            run = line.add_run("_" * 80)
            run.font.name = FONT; run.font.size = Pt(10.5)
        for opt in opts:
            op = doc.add_paragraph(style="List Paragraph")
            op.paragraph_format.left_indent = Cm(0.7)
            run = op.add_run(f"○  {opt}")
            run.font.name = FONT; run.font.size = Pt(10.5)

        # Skip logic
        for cond in (q.condiciones or []):
            skip_p = doc.add_paragraph()
            skip_p.paragraph_format.left_indent = Cm(0.7)
            dest = f"ir a P{cond['ir_a_orden'] + 1}" if cond.get("ir_a_orden") is not None else "Fin de encuesta"
            run = skip_p.add_run(f'↳  Si responde "{cond["si_respuesta"]}" → {dest}')
            run.font.name = FONT; run.font.size = Pt(9)
            run.font.italic = True; run.font.color.rgb = BLUE2

        doc.add_paragraph()  # espacio entre preguntas

    # ── Consideraciones finales (fijas Andersen) ──────────────────
    doc.add_page_break()
    h = doc.add_paragraph()
    r = h.add_run("CONSIDERACIONES FINALES")
    r.font.name = FONT; r.font.size = Pt(11)
    r.font.bold = True; r.font.color.rgb = NAVY
    add_h1_border(h)
    doc.add_paragraph()
    for txt in [
        "La información contenida en este documento es de carácter confidencial y está destinada exclusivamente a su destinatario.",
        "Este documento no implica ningún compromiso de adjudicación por parte de BRAINTRUST SL ni de ANDERSEN CONSULTING.",
        "Los datos, análisis y metodologías aquí descritos son propiedad intelectual de BRAINTRUST SL.",
    ]:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        run = p.add_run(f"• {txt}")
        run.font.name = FONT; run.font.size = Pt(10.5)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def export_survey_pptx(db: Session, survey: Survey) -> bytes:
    """Genera un PPTX de resultados con identidad corporativa Andersen Consulting.
    Layout: slide portada (dark navy) + una slide por pregunta con barras horizontales."""
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Paleta Andersen
    NAVY    = RGBColor(0x17, 0x2D, 0x42)
    WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    BLUE_A  = RGBColor(0x15, 0x81, 0xCE)  # acento
    BLUE_M  = RGBColor(0x21, 0x69, 0x93)  # medio
    GREY_L  = RGBColor(0xE8, 0xE8, 0xE8)
    MUTED_D = RGBColor(0x9A, 0xB0, 0xC6)  # texto secundario dark
    MUTED_L = RGBColor(0x77, 0x77, 0x77)  # texto secundario light
    FONT    = "Arial"

    # Dimensiones 16:9 Andersen (EMU)
    W, H = 12192000, 6858000
    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)
    blank = prs.slide_layouts[6]  # blank

    def rect(slide, x, y, w, h, fill=None, line=False):
        s = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(h))
        if fill:
            s.fill.solid(); s.fill.fore_color.rgb = fill
        else:
            s.fill.background()
        if not line:
            s.line.fill.background()
        else:
            s.line.color.rgb = fill or WHITE
        return s

    def text(slide, txt, x, y, w, h, size=14, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
        tf = tb.text_frame; tf.word_wrap = wrap
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.name = FONT; r.font.size = Pt(size)
        r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color
        return tb

    def footer(slide, dark=True):
        c = MUTED_D if dark else MUTED_L
        text(slide, "ANDERSENCONSULTING", 350000, H - 340000, 4000000, 280000,
             size=8, color=c)

    def left_bar(slide, dark=True):
        rect(slide, 0, 0, 380000, H,
             fill=BLUE_M if dark else GREY_L)
        text(slide, "RESULTS", 50000, H // 2 - 300000, 280000, 600000,
             size=8, bold=True, color=BLUE_A if dark else BLUE_M, align=PP_ALIGN.CENTER)

    total_resp = db.query(SurveyResponse).filter(SurveyResponse.survey_id == survey.id).count()

    # ── PORTADA (dark) ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    rect(sl, 0, 0, W, H, fill=NAVY)
    rect(sl, 0, H - 28000, W, 28000, fill=BLUE_A)   # línea inferior acento
    text(sl, survey.nombre.upper(), 800000, H // 2 - 1000000, W - 1200000, 900000,
         size=34, bold=True, color=WHITE)
    if survey.tema:
        text(sl, survey.tema, 800000, H // 2 - 50000, W - 1200000, 500000,
             size=18, color=BLUE_A)
    text(sl, f"n = {total_resp} respuestas", 800000, H // 2 + 600000, 4000000, 320000,
         size=13, color=MUTED_D)
    footer(sl, dark=True)

    # ── SLIDES POR PREGUNTA ────────────────────────────────────────
    results = compute_results(db, survey, None)

    for qi, qr in enumerate(results["preguntas"]):
        dark = (qi // 3) % 2 == 0  # alterna dark/light cada 3 preguntas
        bg   = NAVY if dark else WHITE
        tc   = WHITE if dark else NAVY
        mut  = MUTED_D if dark else MUTED_L
        bar_bg = RGBColor(0x22, 0x3A, 0x52) if dark else GREY_L

        sl = prs.slides.add_slide(blank)
        rect(sl, 0, 0, W, H, fill=bg)
        left_bar(sl, dark=dark)

        # Número de pregunta + tipo
        text(sl, f"P{qi + 1}", 480000, 300000, 600000, 380000,
             size=28, bold=True, color=BLUE_A, align=PP_ALIGN.LEFT)
        type_lbl = {"single":"Opción única","multiple":"Opción múltiple","yesno":"Sí/No",
                    "likert":"Likert 1–5","nps":"NPS 0–10","abierta":"Abierta"}.get(qr["tipo"], qr["tipo"])
        text(sl, type_lbl, 1050000, 350000, 2500000, 280000,
             size=10, color=BLUE_A if dark else BLUE_M)

        # Texto pregunta
        text(sl, qr["texto"], 480000, 700000, W - 900000, 750000,
             size=17, bold=True, color=tc, wrap=True)

        # n=
        text(sl, f"n = {qr['n']}", 480000, 1400000, 2000000, 260000,
             size=11, color=mut)

        # Media / NPS
        y_info = 1400000
        if qr.get("media") is not None:
            text(sl, f"Media: {qr['media']}", W - 2600000, y_info, 2200000, 260000,
                 size=12, bold=True, color=BLUE_A, align=PP_ALIGN.RIGHT)
        if qr.get("nps") is not None:
            text(sl, f"NPS: {qr['nps']}", W - 2600000, y_info, 2200000, 260000,
                 size=12, bold=True, color=BLUE_A, align=PP_ALIGN.RIGHT)

        if qr["tipo"] == "abierta":
            y = 1800000
            for tv in (qr.get("textos") or [])[:5]:
                card_h = 540000
                rect(sl, 480000, y, W - 960000, card_h,
                     fill=RGBColor(0x1E, 0x35, 0x4D) if dark else RGBColor(0xF0, 0xF4, 0xF8))
                snippet = f"“{tv[:130]}…”" if len(tv) > 130 else f"“{tv}”"
                text(sl, snippet, 620000, y + 100000, W - 1200000, card_h - 200000,
                     size=12, italic=True, color=tc, wrap=True)
                y += card_h + 80000
        else:
            distribs = qr.get("distribucion", [])
            max_n = max((d["n"] for d in distribs), default=1) or 1
            BAR_X  = 2950000   # inicio de barra
            BAR_AW = 7800000   # ancho área de barras
            BAR_H  = 330000
            GAP    = 140000
            y = 1800000
            for d in distribs:
                lbl = d["opcion"]
                if len(lbl) > 26: lbl = lbl[:25] + "…"
                text(sl, lbl, 480000, y + 30000, 2400000, BAR_H,
                     size=11, color=tc)
                # fondo
                rect(sl, BAR_X, y + 30000, BAR_AW, BAR_H - 60000, fill=bar_bg)
                # barra rellena
                bw = int(BAR_AW * d["n"] / max_n) if d["n"] > 0 else 0
                if bw > 0:
                    rect(sl, BAR_X, y + 30000, bw, BAR_H - 60000, fill=BLUE_A)
                # valor
                text(sl, f"{d['n']} · {d['pct']}%",
                     BAR_X + BAR_AW + 80000, y + 30000, 1200000, BAR_H,
                     size=11, color=mut)
                y += BAR_H + GAP

        footer(sl, dark=dark)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
